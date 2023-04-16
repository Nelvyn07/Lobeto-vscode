# -*- coding: utf-8 -*-
"""
This code is used to create .rep file for WAC fails reporting.
The code works for one or multiple child lots.
The WAC fails csv files (one csv file for each child lot) are obtained from LOBETO: FAILS ONLY->EXPORT XLS
The csv files have their default names starting by 'table.csv'. Do not change the name, just place them in the waferworkspace folder, i.e \\vdrsfile5\wafersworkspace$\22FDSOI\Product\Lot\
Lobeto: http://t1onlinev//lobeto3/index.php?mod=disposition&op=details&family_lot_id=URBY03001.000&lot_id=URBY03001.000&operation=FINA-FWET.01&insertion=FWET
A splitsheet from EASI exported as 'DCUBE' format is also needed. This is normally exported automatically from EASI to this folder: \\drsfile5\wafersworkspace$\_automation\EASIsplitsD3\
"""
import tkinter as tk
import glob
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from openpyxl import load_workbook
import webbrowser
import time

# Define the scaling factor
SCALE_FACTOR = 1

def lobeto_link():
    working_folder = folder_path_entry.get()
    shortLot = working_folder.split("\\")[-1]
    url = 'http://t1onlinev//lobeto3/index.php?mod=disposition&op=details&family_lot_id='+shortLot+'.000&lot_id='+shortLot+'.000&operation=FINA-FWET.01&insertion=FWET'
    time.sleep(1)
    webbrowser.open(url)

def rep_creator():
    working_folder = folder_path_entry.get()
    working_folder = working_folder+"\\"
    shortLot = working_folder.split("\\")[-2]

    splitsheet_folder = '\\\\vdrsfile5\\wafersworkspace$\\_automation\\EASIsplitsD3\\'

    # Set the path to the folder containing the CSV files
    template_path = '\\\\vdrsfile5\\wafersworkspace$\\22FDSOI\\Definition_Corners\\wac_fails_auto_report\\' 
    

    # List all CSV files in the current directory that start with 'lot'
    wac_csv_files = glob.glob(working_folder+'table*.csv')

    # Create an empty list to store the dataframes
    df_list = []

    # Loop through each CSV file and read it into a dataframe, then append it to the list. df_list is a list of dataframes.
    for filename in wac_csv_files:
        #df = pd.read_csv(filename, error_bad_lines=False)
        df = pd.read_csv(filename, on_bad_lines='skip')
        #df = df[~df['FLOW'].str.contains('Pass')] #removed the rows containing Pass.
        df_list.append(df)

    # Concatenate all the dataframes into a single one
    try:
        df = pd.concat(df_list)
        WAC_fails_file_label.config(text="WAC fails table found in the working folder.", font=("TkDefaultFont", int(10*SCALE_FACTOR)), foreground="green")
    except(ValueError):
        WAC_fails_file_label.config(text="Folder path not correct or WAC fails table not found. Get WAC fails from Lobeto.", font=("TkDefaultFont", int(10*SCALE_FACTOR)), foreground="red")
        return
    
    df = df[~df['FLOW'].str.contains('Pass')] #removed the rows containing Pass.

    #delete all Unnamed Columns in a single code of line using regex
    df.drop(df.filter(regex="Unnamed"),axis=1, inplace=True)

    # Define a list of columns to be dropped from df
    to_drop = ["FLOW", "PMS", "SCRIBE", "SLOT", "STTI", "INSERTION", "TS", "TP", "PC", "PS", "SL", "LS", "TOTAL", "PASS/FAIL"]

    # Drop the specified columns from df
    df.drop(columns = to_drop, inplace = True)

    # Rename the 'ID' column to 'WAFER_NUMBER' so that df and dcube share the same column name and can be merged.
    df = df.rename(columns={'ID': 'WAFER_NUMBER'})
    df['WAFER_NUMBER'] = df['WAFER_NUMBER'].apply(lambda x: int(x)) #WAFER_NUM in df and dcube has to be in the same format: int works.

    try:
        Dcube_Split = glob.glob(splitsheet_folder+'Dcube_Split_' + shortLot + '*.csv')[0]
        splitfile_label.config(text="Splitfile found in " + splitsheet_folder, font=("TkDefaultFont", int(10*SCALE_FACTOR)), foreground="green")       
    except(IndexError):    
        splitfile_label.config(text="No splitfile found.\nFill the EASI splitsheet, click on the  " + splitsheet_folder, font=("TkDefaultFont", int(10*SCALE_FACTOR)), foreground="red")
        return

    dcube = pd.read_csv(Dcube_Split) 

    # extract the string between parentheses and join the strings for the wafers which are composed of more than one split information, e.g FF, APMOM
    dcube = dcube.groupby("WAFER_NUMBER")["SPLIT"].apply(lambda x: ", ".join(x.astype(str))).reset_index()
    dcube['WAFER_NUMBER'] = dcube['WAFER_NUMBER'].apply(lambda x: int(x))
                                       
    df_merged = pd.merge(df, dcube, on= 'WAFER_NUMBER')

    # Set the 'WAFER_NUMBER' column as the index of df_merged and sort the rows by index
    df_merged = df_merged.set_index('WAFER_NUMBER')
    df_merged = df_merged.sort_index()

    # Sort the columns of df_merged alphabetically
    df_merged = df_merged.sort_index(axis=1)

    # Reset the index of df_merged
    df_merged = df_merged.reset_index()

    # Reorder the columns of df_merged so that the 'SPLIT' column is the first column
    df_merged.insert(0,'SPLIT',df_merged.pop('SPLIT'))

    # Define a function to fill cells with the string 'WAC fail' if the value is less than 50
    def wac_fill(x):
        x = float(x) #there are some strings which appear in the table creating an error 
        if x<50:
            return 'WAC fail'
        else:
            return '-'
        
    # Apply the wac_fill function to all columns after WAFER_NUMBER
    df_merged.iloc[:,2:]=df_merged.iloc[:,2:].applymap(wac_fill)

    df_merged.to_excel(working_folder+'wac_fails_summary.xlsx', index=False)


    #-----------section to create .plo for plotting the wac fails------------------

    try:
        # param.plo is used as template to create the wac_fails.plo
        df_plo = pd.read_csv(template_path+'wac_fails.plo.csv')
        template_label.config(text="")
    except(FileNotFoundError):
        template_label.config(text="Template files not found. Check that the path "+ template_path + " has not been moved or renamed.\n It is recommended not to change the .pot template as issues with the auto comments could happen." , font=("TkDefaultFont", int(10*SCALE_FACTOR)), foreground="red")
        return

    num_param = len(df_merged.columns)-2

    # repeat the first row n times
    df_plo = pd.concat([df_plo]*num_param, ignore_index=True)

    # 'name' is integer incrementing from 1.
    df_plo['name'] = range(1, len(df_plo) + 1)

    #'title': parameter name 
    df_plo.loc[:, 'title'] = list(df_merged.columns[2:])
    df_plo['y'] = 'swet_'+ df_plo['title']

    df_plo['abslog'] = df_plo['y'].apply(lambda x: 'logy' if 'IOE' in x else 'absy')

    df_plo.to_csv(working_folder+ shortLot +"_wac_fails.plo.csv", index=False)


    #-------------------section to create .pot-----------------------

    # Open the template PowerPoint presentation
    prs = Presentation(template_path+'wac_fails.pot.pptx')

    # Access slide 1 and add shortLot
    slide = prs.slides[0]
    slide.shapes[1].text = shortLot + ' @FINA-FWET.01\nTEQ DRS TD'

    # Access slide 2
    slide = prs.slides[1]

    workbook = load_workbook(working_folder+'wac_fails_summary.xlsx')
    worksheet = workbook.active


    # Insert the Excel table into the slide
    table = slide.shapes.add_table(rows=worksheet.max_row, cols=worksheet.max_column, left=Inches(0.3), top=Inches(1), width=Inches(12), height=Inches(2)).table

    for i in range(worksheet.max_row):
        for j in range(worksheet.max_column):
            table.cell(i, j).text = str(worksheet.cell(i + 1, j + 1).value)

    for cell in table.iter_cells():
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(10)

    # Duplicate the 3rd slide (template for plots) n times and add the corresponding title and comment for each slide.
    def duplicate_slide(pres, index):
        template = pres.slides[index]

        blank_slide_layout = pres.slide_layouts[23]
        # blank slide corresponds to the 24th GF slide layouts. As layouts start with index 0, we use here 23. 
        # Check that the layout index did not change when we have a new powerpoint template!
            
        copied_slide = pres.slides.add_slide(blank_slide_layout)

        for shp in template.shapes:
            el = shp.element
            import copy
            newel = copy.deepcopy(el)
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

        return copied_slide


    for index in range(len(df_merged.columns[2:])):
        new_slide = duplicate_slide(prs,2)
        new_slide.shapes[2].text_frame.text = str(index+1)
        from pptx.enum.text import PP_ALIGN
        new_slide.shapes[2].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        new_slide.shapes[0].text_frame.text = 'WAC fail at FWET: ' + str(df_merged.columns[index+2])
        title_para = new_slide.shapes[0].text_frame.paragraphs[0]
        title_para.font.name = "Arial Black"
        title_para.font.size = Pt(28)
        title_para.font.color.rgb = RGBColor(0x4B, 0x00, 0x82)
        
        param = df_merged.columns[index+2]
        df1 = df_merged.loc[df_merged[param]=='WAC fail']
        df1 = df1.groupby("SPLIT", group_keys=True)["WAFER_NUMBER"].apply(lambda x: "&".join(x.astype(str)))

        fails = dict(df1)
        line =[]
        for key, value in fails.items():
                words = str(value) + " (" + str(key) + ")"
                line.append(words)
                        
        line = ", ".join(line)
        line = "WAC fail due to cornering on wafer alias:\n " + line +"\n No issue for shipment."
        
        new_slide.shapes[1].text_frame.text = line
        
        for para in new_slide.shapes[1].text_frame.paragraphs:
            para.font.name = "Arial"
            para.font.size = Pt(18) 
    
    # to remove the empty template slide    
    xml_slides = prs.slides._sldIdLst  
    slides = list(xml_slides)
    xml_slides.remove(slides[2])


    # Save the PowerPoint presentation
    prs.save(working_folder+ shortLot +'_wac_fails.pot.pptx')


    #-------------------section to create .rep-----------------------
    df_rep = pd.read_csv(template_path+'wac_fails.rep.csv')
    df_rep.iloc[2,2] = working_folder+ shortLot +'_wac_fails.pot.pptx'
    df_rep.iloc[3,2] = working_folder+ shortLot +'_wac_fails.plo.csv'
    df_rep.iloc[3,4] = working_folder+ shortLot +'_wac_fails.pptx'

    #spl_path = glob.glob(working_folder+'*.spl.csv')[0]
    spl_path = working_folder+shortLot+'_FINAFWETFWET_AUTO.SPL.CSV'
    df_rep.iloc[0,2] = spl_path

    config_folder = working_folder.split(shortLot+"\\")[0] + '_Config\\'
    lim_path = glob.glob(config_folder + '*.lim.csv')[0]
    df_rep.iloc[1,2] = lim_path

    df_rep.to_csv(working_folder+ shortLot +'_wac_fails.rep.csv', index=False)

    rep_label.config(text="The .rep file has been created. \nFor MPW lots, you might need to change the limit file in the .rep to match to the correct Product ID.\nYou might need to modify the splitsheet if the report is for subset of wafers in a child lot.", font=("TkDefaultFont", int(10*SCALE_FACTOR)), foreground="green")
        

# Create the main window
window = tk.Tk()
window.title("rep file creator for WAC fails at FWET on 22FDX corner lots")

# Create the folder_path input field
folder_path_label = tk.Label(window, text="Enter the working folder path in the format \\\\vdrsfile5\wafersworkspace$\\22FDSOI\Product\Lot : ", font=("TkDefaultFont", int(10*SCALE_FACTOR)))

folder_path_entry = tk.Entry(window, font=("TkDefaultFont", int(10*SCALE_FACTOR)), width=int(100*SCALE_FACTOR))

lobeto_label = tk.Label(window, text="Get the WAC fails csv files (one csv file for each child lot) from Lobeto: click on FAILS ONLY->EXPORT XLS). \nThe csv files have their default names starting by 'table.csv'. Do not change the name, just place them in the waferworkspace folder.\n If you have to plot the WAC fails on the mother lot (.000) only, you can directly get the link to Lobeto by clicking on the Lobeto link.", font=("TkDefaultFont", int(10*SCALE_FACTOR)), foreground="blue")

# Create the button
lobeto_button = tk.Button(window, text="Lobeto website link \nto mother lot", font=("TkDefaultFont", int(10*SCALE_FACTOR)), command=lobeto_link)

WAC_fails_file_label = tk.Label(window, text="", font=("TkDefaultFont", int(10*SCALE_FACTOR)))

splitfile_label = tk.Label(window, text="", font=("TkDefaultFont", int(10*SCALE_FACTOR)))

template_label = tk.Label(window, text="", font=("TkDefaultFont", int(10*SCALE_FACTOR)))

# Create the button
rep_button = tk.Button(window, text="Create .rep", font=("TkDefaultFont", int(10*SCALE_FACTOR)), command=rep_creator)

# Create the rep output message
rep_label = tk.Label(window, text="", font=("TkDefaultFont", int(10*SCALE_FACTOR)))


folder_path_label.grid(row=0, column=0)
folder_path_entry.grid(row=1, column=0, pady=10*SCALE_FACTOR)
lobeto_label.grid(row=2, column=0, pady=10*SCALE_FACTOR)
lobeto_button.grid(row=2, column=1, padx=10)
rep_button.grid(row=3, column=0, pady=10*SCALE_FACTOR)
WAC_fails_file_label.grid(row=4, column=0, pady=10*SCALE_FACTOR)
splitfile_label.grid(row=5, column=0, pady=10*SCALE_FACTOR)
template_label.grid(row=6, column=0, pady=10*SCALE_FACTOR)
rep_label.grid(row=6, column=0, pady=10*SCALE_FACTOR)

# Scale up the window dimensions
window_width = int(1000*SCALE_FACTOR)
window_height = int(400*SCALE_FACTOR)
window.geometry(f"{window_width}x{window_height}")

# Run the main loop
window.mainloop()