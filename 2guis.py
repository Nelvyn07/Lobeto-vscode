import tkinter as tk
import webbrowser
import time
import glob
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from openpyxl import load_workbook
import requests
import xml.etree.ElementTree as ET

SCALE_FACTOR = 1

class MainWindow:
    def __init__(self, master):
        self.master = master
        master.title("Tasks Automation for 22FDX corner lots")

        self.first_button = tk.Button(master, text="WAC fails", font=("TkDefaultFont", int(20*SCALE_FACTOR)), command=self.open_first_window)
        self.first_button.pack()

        self.second_button = tk.Button(master, text="XML for ERFs", font=("TkDefaultFont", int(20*SCALE_FACTOR)), command=self.open_second_window)
        self.second_button.pack()
        # Scale up the window dimensions
        self.master_width = int(400*SCALE_FACTOR)
        self.master_height = int(200*SCALE_FACTOR)
        self.master.geometry(f"{self.master_width}x{self.master_height}")

    def open_first_window(self):
        self.first_window = tk.Toplevel(self.master)
        self.first_window.title("rep file creator for WAC fails at FWET on 22FDX corner lots")
        # Create the folder_path input field
        self.folder_path_label = tk.Label(self.first_window, text="Enter the working folder path in the format \\\\vdrsfile5\wafersworkspace$\\22FDSOI\Product\Lot : ", font=("TkDefaultFont", int(10*SCALE_FACTOR)))

        self.folder_path_entry = tk.Entry(self.first_window, font=("TkDefaultFont", int(10*SCALE_FACTOR)), width=int(100*SCALE_FACTOR))

        self.lobeto_label = tk.Label(self.first_window, text="Get the WAC fails csv files (one csv file for each child lot) from Lobeto: click on FAILS ONLY->EXPORT XLS). \nThe csv files have their default names starting by 'table.csv'. Do not change the name, just place them in the waferworkspace folder.\n If you have to plot the WAC fails on the mother lot (.000) only, you can directly get the link to Lobeto by clicking on the Lobeto link.", font=("TkDefaultFont", int(10*SCALE_FACTOR)), foreground="blue")

        # Create the button
        self.lobeto_button = tk.Button(self.first_window, text="Lobeto website link \nto mother lot", font=("TkDefaultFont", int(10*SCALE_FACTOR)), command=self.lobeto_link)

        self.WAC_fails_file_label = tk.Label(self.first_window, text="", font=("TkDefaultFont", int(10*SCALE_FACTOR)))

        self.splitfile_label = tk.Label(self.first_window, text="", font=("TkDefaultFont", int(10*SCALE_FACTOR)))

        self.template_label = tk.Label(self.first_window, text="", font=("TkDefaultFont", int(10*SCALE_FACTOR)))

        # Create the button
        self.rep_button = tk.Button(self.first_window, text="Create .rep", font=("TkDefaultFont", int(10*SCALE_FACTOR)), command=self.rep_creator)

        # Create the rep output message
        self.rep_label = tk.Label(self.first_window, text="", font=("TkDefaultFont", int(10*SCALE_FACTOR)))
       

        self.folder_path_label.grid(row=0, column=0)
        self.folder_path_entry.grid(row=1, column=0, pady=10*SCALE_FACTOR)
        self.lobeto_label.grid(row=2, column=0, pady=10*SCALE_FACTOR)
        self.lobeto_button.grid(row=2, column=1, padx=10)
        self.rep_button.grid(row=3, column=0, pady=10*SCALE_FACTOR)
        self.WAC_fails_file_label.grid(row=4, column=0, pady=10*SCALE_FACTOR)
        self.splitfile_label.grid(row=5, column=0, pady=10*SCALE_FACTOR)
        self.template_label.grid(row=6, column=0, pady=10*SCALE_FACTOR)
        self.rep_label.grid(row=6, column=0, pady=10*SCALE_FACTOR)

        # Scale up the window dimensions
        self.first_window_width = int(1000*SCALE_FACTOR)
        self.first_window_height = int(400*SCALE_FACTOR)
        self.first_window.geometry(f"{self.first_window_width}x{self.first_window_height}")

    def open_second_window(self):
        self.second_window = tk.Toplevel(self.master)
        self.second_window.title("EASI to XML for 22FDX corner ERFs")

        # Create the user id input field
        self.user_id_label = tk.Label(self.second_window, text="User ID:", font=("TkDefaultFont", int(10*SCALE_FACTOR)))
        self.user_id_label.pack()
        self.user_id_entry = tk.Entry(self.second_window, font=("TkDefaultFont", int(10*SCALE_FACTOR)), width=int(20*SCALE_FACTOR))
        self.user_id_entry.insert(0, "yandee")
        self.user_id_entry.pack()

        # Create the password input field
        self.password_label = tk.Label(self.second_window, text="Password:", font=("TkDefaultFont", int(10*SCALE_FACTOR)))
        self.password_label.pack()
        self.password_entry = tk.Entry(self.second_window, font=("TkDefaultFont", int(10*SCALE_FACTOR)), width=int(20*SCALE_FACTOR), show="*")
        self.password_entry.insert(0, "Homeoffice30")
        self.password_entry.pack()

        # Create the show password checkbutton
        self.show_password = tk.BooleanVar()
        self.show_password.set(False)
        self.show_password_checkbutton = tk.Checkbutton(self.second_window, text="Show password", variable=self.show_password, font=("TkDefaultFont", int(10*SCALE_FACTOR)), command=self.toggle_password_visibility)
        self.show_password_checkbutton.pack()


        # Create the correct/incorrect password label
        self.password_validity_label = tk.Label(self.second_window, text="", font=("TkDefaultFont", int(10*SCALE_FACTOR)))
        self.password_validity_label.pack()

        # Create the erf_id input field
        self.erf_id_label = tk.Label(self.second_window, text="Enter the EASI ErfID:", font=("TkDefaultFont", int(10*SCALE_FACTOR)))
        self.erf_id_label.pack()
        self.erf_id_entry = tk.Entry(self.second_window, font=("TkDefaultFont", int(10*SCALE_FACTOR)), width=int(20*SCALE_FACTOR))
        self.erf_id_entry.pack()

        # Create the button
        self.XML_button = tk.Button(self.second_window, text="XML", font=("TkDefaultFont", int(10*SCALE_FACTOR)), command=self.easi_to_xml)
        self.XML_button.pack()

        # Create the XML label
        self.XML_label = tk.Label(self.second_window, text="", font=("TkDefaultFont", int(10*SCALE_FACTOR)))
        self.XML_label.pack()

        # Scale up the window dimensions
        self.second_window_width = int(400*SCALE_FACTOR)
        self.second_window_height = int(300*SCALE_FACTOR)
        self.second_window.geometry(f"{self.second_window_width}x{self.second_window_height}")


    def lobeto_link(self):
        working_folder = self.folder_path_entry.get()
        shortLot = working_folder.split("\\")[-1]
        url = 'http://t1onlinev//lobeto3/index.php?mod=disposition&op=details&family_lot_id='+shortLot+'.000&lot_id='+shortLot+'.000&operation=FINA-FWET.01&insertion=FWET'
        time.sleep(1)
        webbrowser.open(url)

    def rep_creator(self):
        working_folder = self.folder_path_entry.get()
        working_folder = working_folder+"\\"
        shortLot = working_folder.split("\\")[-2]

        splitsheet_folder = '\\\\vdrsfile5\\wafersworkspace$\\_automation\\EASIsplitsD3\\'

        # Set the path to the folder containing the CSV files
        template_path = '\\\\vdrsfile5\\wafersworkspace$\\22FDSOI\\Definition_Corners\\wac_fails_auto_report\\' 
        

        # List all CSV files in the current directory that start with 'lot'
        wac_csv_files = glob.glob(working_folder+'table*.csv')

        # Create an empty list to store the dataframes
        df_list = []

        # Loop through each CSV file and read it into a dataframe, then append it to the list. df_list is a list of dataframes.
        for filename in wac_csv_files:
            #df = pd.read_csv(filename, error_bad_lines=False)
            df = pd.read_csv(filename, on_bad_lines='skip')
            #df = df[~df['FLOW'].str.contains('Pass')] #removed the rows containing Pass.
            df_list.append(df)

        # Concatenate all the dataframes into a single one
        try:
            df = pd.concat(df_list)
            self.WAC_fails_file_label.config(text="WAC fails table found in the working folder.", font=("TkDefaultFont", int(10*SCALE_FACTOR)), foreground="green")
        except(ValueError):
            self.WAC_fails_file_label.config(text="Folder path not correct or WAC fails table not found. Get WAC fails from Lobeto.", font=("TkDefaultFont", int(10*SCALE_FACTOR)), foreground="red")
            return
        
        df = df[~df['FLOW'].str.contains('Pass')] #removed the rows containing Pass.

        #delete all Unnamed Columns in a single code of line using regex
        df.drop(df.filter(regex="Unnamed"),axis=1, inplace=True)

        # Define a list of columns to be dropped from df
        to_drop = ["FLOW", "PMS", "SCRIBE", "SLOT", "STTI", "INSERTION", "TS", "TP", "PC", "PS", "SL", "LS", "TOTAL", "PASS/FAIL"]

        # Drop the specified columns from df
        df.drop(columns = to_drop, inplace = True)

        # Rename the 'ID' column to 'WAFER_NUMBER' so that df and dcube share the same column name and can be merged.
        df = df.rename(columns={'ID': 'WAFER_NUMBER'})
        df['WAFER_NUMBER'] = df['WAFER_NUMBER'].apply(lambda x: int(x)) #WAFER_NUM in df and dcube has to be in the same format: int works.

        try:
            Dcube_Split = glob.glob(splitsheet_folder+'Dcube_Split_' + shortLot + '*.csv')[0]
            self.splitfile_label.config(text="Splitfile found in " + splitsheet_folder, font=("TkDefaultFont", int(10*SCALE_FACTOR)), foreground="green")       
        except(IndexError):    
            self.splitfile_label.config(text="No splitfile found.\nFill the EASI splitsheet, click on the  " + splitsheet_folder, font=("TkDefaultFont", int(10*SCALE_FACTOR)), foreground="red")
            return

        dcube = pd.read_csv(Dcube_Split) 

        # extract the string between parentheses and join the strings for the wafers which are composed of more than one split information, e.g FF, APMOM
        dcube = dcube.groupby("WAFER_NUMBER")["SPLIT"].apply(lambda x: ", ".join(x.astype(str))).reset_index()
        dcube['WAFER_NUMBER'] = dcube['WAFER_NUMBER'].apply(lambda x: int(x))
                                        
        df_merged = pd.merge(df, dcube, on= 'WAFER_NUMBER')

        # Set the 'WAFER_NUMBER' column as the index of df_merged and sort the rows by index
        df_merged = df_merged.set_index('WAFER_NUMBER')
        df_merged = df_merged.sort_index()

        # Sort the columns of df_merged alphabetically
        df_merged = df_merged.sort_index(axis=1)

        # Reset the index of df_merged
        df_merged = df_merged.reset_index()

        # Reorder the columns of df_merged so that the 'SPLIT' column is the first column
        df_merged.insert(0,'SPLIT',df_merged.pop('SPLIT'))

        # Define a function to fill cells with the string 'WAC fail' if the value is less than 50
        def wac_fill(x):
            x = float(x) #there are some strings which appear in the table creating an error 
            if x<50:
                return 'WAC fail'
            else:
                return '-'
            
        # Apply the wac_fill function to all columns after WAFER_NUMBER
        df_merged.iloc[:,2:]=df_merged.iloc[:,2:].applymap(wac_fill)

        df_merged.to_excel(working_folder+'wac_fails_summary.xlsx', index=False)


        #-----------section to create .plo for plotting the wac fails------------------

        try:
            # param.plo is used as template to create the wac_fails.plo
            df_plo = pd.read_csv(template_path+'wac_fails.plo.csv')
            self.template_label.config(text="")
        except(FileNotFoundError):
            self.template_label.config(text="Template files not found. Check that the path "+ template_path + " has not been moved or renamed.\n It is recommended not to change the .pot template as issues with the auto comments could happen." , font=("TkDefaultFont", int(10*SCALE_FACTOR)), foreground="red")
            return

        num_param = len(df_merged.columns)-2

        # repeat the first row n times
        df_plo = pd.concat([df_plo]*num_param, ignore_index=True)

        # 'name' is integer incrementing from 1.
        df_plo['name'] = range(1, len(df_plo) + 1)

        #'title': parameter name 
        df_plo.loc[:, 'title'] = list(df_merged.columns[2:])
        df_plo['y'] = 'swet_'+ df_plo['title']

        df_plo['abslog'] = df_plo['y'].apply(lambda x: 'logy' if 'IOE' in x else 'absy')

        df_plo.to_csv(working_folder+ shortLot +"_wac_fails.plo.csv", index=False)


        #-------------------section to create .pot-----------------------

        # Open the template PowerPoint presentation
        prs = Presentation(template_path+'wac_fails.pot.pptx')

        # Access slide 1 and add shortLot
        slide = prs.slides[0]
        slide.shapes[1].text = shortLot + ' @FINA-FWET.01\nTEQ DRS TD'

        # Access slide 2
        slide = prs.slides[1]

        workbook = load_workbook(working_folder+'wac_fails_summary.xlsx')
        worksheet = workbook.active


        # Insert the Excel table into the slide
        table = slide.shapes.add_table(rows=worksheet.max_row, cols=worksheet.max_column, left=Inches(0.3), top=Inches(1), width=Inches(12), height=Inches(2)).table

        for i in range(worksheet.max_row):
            for j in range(worksheet.max_column):
                table.cell(i, j).text = str(worksheet.cell(i + 1, j + 1).value)

        for cell in table.iter_cells():
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)

        # Duplicate the 3rd slide (template for plots) n times and add the corresponding title and comment for each slide.
        def duplicate_slide(pres, index):
            template = pres.slides[index]

            blank_slide_layout = pres.slide_layouts[23]
            # blank slide corresponds to the 24th GF slide layouts. As layouts start with index 0, we use here 23. 
            # Check that the layout index did not change when we have a new powerpoint template!
                
            copied_slide = pres.slides.add_slide(blank_slide_layout)

            for shp in template.shapes:
                el = shp.element
                import copy
                newel = copy.deepcopy(el)
                copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

            return copied_slide


        for index in range(len(df_merged.columns[2:])):
            new_slide = duplicate_slide(prs,2)
            new_slide.shapes[2].text_frame.text = str(index+1)
            from pptx.enum.text import PP_ALIGN
            new_slide.shapes[2].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            new_slide.shapes[0].text_frame.text = 'WAC fail at FWET: ' + str(df_merged.columns[index+2])
            title_para = new_slide.shapes[0].text_frame.paragraphs[0]
            title_para.font.name = "Arial Black"
            title_para.font.size = Pt(28)
            title_para.font.color.rgb = RGBColor(0x4B, 0x00, 0x82)
            
            param = df_merged.columns[index+2]
            df1 = df_merged.loc[df_merged[param]=='WAC fail']
            df1 = df1.groupby("SPLIT", group_keys=True)["WAFER_NUMBER"].apply(lambda x: "&".join(x.astype(str)))

            fails = dict(df1)
            line =[]
            for key, value in fails.items():
                    words = str(value) + " (" + str(key) + ")"
                    line.append(words)
                            
            line = ", ".join(line)
            line = "WAC fail due to cornering on wafer alias:\n " + line +"\n No issue for shipment."
            
            new_slide.shapes[1].text_frame.text = line
            
            for para in new_slide.shapes[1].text_frame.paragraphs:
                para.font.name = "Arial"
                para.font.size = Pt(18) 
        
        # to remove the empty template slide    
        xml_slides = prs.slides._sldIdLst  
        slides = list(xml_slides)
        xml_slides.remove(slides[2])


        # Save the PowerPoint presentation
        prs.save(working_folder+ shortLot +'_wac_fails.pot.pptx')


        #-------------------section to create .rep-----------------------
        df_rep = pd.read_csv(template_path+'wac_fails.rep.csv')
        df_rep.iloc[2,2] = working_folder+ shortLot +'_wac_fails.pot.pptx'
        df_rep.iloc[3,2] = working_folder+ shortLot +'_wac_fails.plo.csv'
        df_rep.iloc[3,4] = working_folder+ shortLot +'_wac_fails.pptx'

        #spl_path = glob.glob(working_folder+'*.spl.csv')[0]
        spl_path = working_folder+shortLot+'_FINAFWETFWET_AUTO.SPL.CSV'
        df_rep.iloc[0,2] = spl_path

        config_folder = working_folder.split(shortLot+"\\")[0] + '_Config\\'
        lim_path = glob.glob(config_folder + '*.lim.csv')[0]
        df_rep.iloc[1,2] = lim_path

        df_rep.to_csv(working_folder+ shortLot +'_wac_fails.rep.csv', index=False)

        self.rep_label.config(text="The .rep file has been created. \nFor MPW lots, you might need to change the limit file in the .rep to match to the correct Product ID.\nYou might need to modify the splitsheet if the report is for subset of wafers in a child lot.", font=("TkDefaultFont", int(10*SCALE_FACTOR)), foreground="green")
            
    def easi_to_xml(self):
        erf_id = self.erf_id_entry.get()

        url = 'http://tecnet.gfoundries.com/easi/getAjax.php?Step=getsplit&erfid='+ erf_id
        # using http because https causes SSLError
        
        proxies = {
        "http": 'http://dewwwp1.gfoundries.com:74',
        "https": 'http://dewwwp1.gfoundries.com:74',
            }

        passwd = self.password_entry.get()
        user_id = self.user_id_entry.get()
        
        
        try:
            r = requests.get(url, proxies=proxies, auth=(user_id, passwd))
            xml_text = r.text
            xml_text = xml_text.split('SPLIT_INFO')[1]
            self.password_validity_label.config(text="Correct password", font=("TkDefaultFont", int(10*SCALE_FACTOR)), foreground="blue")
        except (IndexError, requests.exceptions.RequestException):
            self.password_validity_label.config(text="Check password!!!", font=("TkDefaultFont", int(10*SCALE_FACTOR)), foreground="red")
            return
        
        xml_text = '<SPLIT_INFO' + xml_text + 'SPLIT_INFO>'
        #Error if the texts before and after SPLIT_INFO are not removed

        folder_path = '\\\\vdrsfile5\\wafersworkspace$\\22FDSOI\\Definition_ya\\XML\\'
        
        with open(folder_path + 'Raw.xml', 'w') as f:
            f.write(xml_text)
        
        tree = ET.parse(folder_path + 'Raw.xml')
        root = tree.getroot()
        
        main_route = root[0].text
        #main route information is needed as it changes from Product to Product and is required for the import in the ERF.
        
        split_info = []
        for instance in root.iter('SPLIT_INSTANCEID'):
            process = instance[0].text.split(sep = '-', maxsplit =1)[0] # to get LW,XW,3PL and SPK
            for spl_grp in instance.iter('SPLIT_GROUP'):
                split = process + ':'+ spl_grp.attrib['splitShort']
                for waf in spl_grp:
                    wafers = split + ':' + waf.text
                    split_info.append(wafers)
        
        df = pd.DataFrame(split_info)
        df.columns = ['Raw']
        #There are some rows which are missing ';' at the end. This causes an issue when we have to concatenate the wafers. 
        #The ~ operator is used to invert the boolean mask so that it selects the rows that do not end with ;
        mask = ~df['Raw'].str.endswith(';') 
        df.loc[mask, 'Raw'] = df.loc[mask, 'Raw'] + ';'
        df['PD'] = df['Raw'].str.split(':', expand =True)[0]
        df['Corner'] = df['Raw'].str.split(':', expand =True)[1]
        df['Wafers'] = df['Raw'].str.split(':', expand =True)[2]
        
        SS_row = df.loc[(df['PD']=='SPIK') & (df['Corner'] =='SS_1.5S')]
        FF_row = df.loc[(df['PD']=='SPIK') & (df['Corner'] =='FF_1.5S')]
        rows_985C = df.loc[(df['PD']=='3PL') & (df['Corner'] !='SS_1.5S') & (df['Corner'] !='FF_1.5S')] 
        
        SS_wafers = SS_row.iloc[0][-1] # Wafer assignment is the last column
        FF_wafers = FF_row.iloc[0][-1]
        other_wafers = ''.join(rows_985C['Wafers']) #concatenate all rows in the 'Wafers' column
        
        tree = ET.parse(folder_path + 'new_3PL+10C\\' +'Template.xml')
        root = tree.getroot()
        
        root[0].text = main_route
        
        for spl_grp in root.iter('SPLIT_GROUP'):
            if spl_grp.attrib['splitShort'] == 'SS_1.5S':
                for waf in spl_grp:
                    waf.text = SS_wafers
                    
            elif spl_grp.attrib['splitShort'] == 'FF_1.5S':
                for waf in spl_grp:
                    waf.text = FF_wafers
        
            elif spl_grp.attrib['splitShort'] == '985C':
                for waf in spl_grp:
                    waf.text = other_wafers
                    
        tree.write(folder_path + 'New.xml')
        
        self.XML_label.config(text="The XML file has been created as New.xml", font=("TkDefaultFont", int(10*SCALE_FACTOR)), foreground="blue")



    def toggle_password_visibility(self):
        if self.show_password.get():
            self.password_entry.config(show="")
        else:
            self.password_entry.config(show="*")

root = tk.Tk()
app = MainWindow(root)
root.mainloop()





